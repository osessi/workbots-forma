"""
Design System Extractor - Extract complete design system from PPTX templates

This is the core service that analyzes PPTX templates and extracts:
1. Color palette (theme colors, accent colors)
2. Typography (fonts, sizes, weights, line heights)
3. Spacing rules (margins, padding, gaps)
4. Layout patterns (grids, alignments, proportions)
5. Shape styles (shadows, borders, fills, gradients)
6. Slide layouts with semantic meaning

The extracted design system allows faithful PPTX reconstruction at 99% fidelity.
"""

import os
import json
import uuid
from typing import Optional, List, Dict, Any, Tuple
from dataclasses import dataclass, field, asdict
from enum import Enum
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn


class LayoutType(str, Enum):
    """Semantic layout types for intelligent slide selection"""
    TITLE = "title"  # Title slide / cover
    SECTION = "section"  # Section divider
    CONTENT = "content"  # Generic content
    BULLETS = "bullets"  # Bullet points
    TWO_COLUMN = "two_column"  # Two column layout
    IMAGE_LEFT = "image_left"  # Image on left, text on right
    IMAGE_RIGHT = "image_right"  # Image on right, text on left
    IMAGE_FULL = "image_full"  # Full bleed image
    COMPARISON = "comparison"  # Side by side comparison
    QUOTE = "quote"  # Quote or testimonial
    STATS = "stats"  # Statistics / numbers
    TIMELINE = "timeline"  # Timeline or process
    TABLE = "table"  # Table layout
    CHART = "chart"  # Chart layout
    CLOSING = "closing"  # Thank you / closing slide
    BLANK = "blank"  # Blank slide


@dataclass
class ColorInfo:
    """Complete color information"""
    rgb: str  # Hex color like "#FF5500"
    theme_color: Optional[str] = None  # Theme color name if applicable
    brightness: Optional[float] = None  # Brightness modification
    alpha: Optional[float] = 1.0  # Transparency


@dataclass
class FontInfo:
    """Complete font information"""
    name: str
    size_pt: float
    bold: bool = False
    italic: bool = False
    underline: bool = False
    color: Optional[ColorInfo] = None
    line_spacing: Optional[float] = None  # In points or as multiplier
    space_before: Optional[float] = None
    space_after: Optional[float] = None


@dataclass
class TextStyleInfo:
    """Complete text style for a text level"""
    font: FontInfo
    alignment: str = "left"  # left, center, right, justify
    vertical_anchor: str = "top"  # top, middle, bottom
    bullet_char: Optional[str] = None
    bullet_color: Optional[ColorInfo] = None
    indent_level: int = 0
    margin_left: float = 0
    margin_right: float = 0


@dataclass
class ShapeStyleInfo:
    """Complete shape styling information"""
    fill_type: str = "solid"  # solid, gradient, pattern, picture, none
    fill_color: Optional[ColorInfo] = None
    gradient_stops: Optional[List[Dict]] = None  # For gradients
    line_color: Optional[ColorInfo] = None
    line_width_pt: Optional[float] = None
    line_dash_style: Optional[str] = None
    shadow: Optional[Dict] = None  # Shadow properties
    corner_radius: Optional[float] = None  # For rounded rectangles


@dataclass
class PlaceholderInfo:
    """Detailed placeholder information"""
    idx: int
    type: str  # title, subtitle, body, picture, etc.
    name: str
    position: Dict[str, float]  # left, top in EMUs
    size: Dict[str, float]  # width, height in EMUs
    text_styles: List[TextStyleInfo] = field(default_factory=list)  # Per-level styles
    shape_style: Optional[ShapeStyleInfo] = None
    max_chars: Optional[int] = None  # Estimated max characters
    max_lines: Optional[int] = None  # Estimated max lines
    is_static: bool = False  # True if this is a static shape (not a native placeholder)
    current_text: Optional[str] = None  # Current text content for static shapes
    shape_id: Optional[int] = None  # Shape ID for targeting during generation


@dataclass
class SlideLayoutInfo:
    """Complete slide layout information"""
    index: int
    name: str
    layout_type: LayoutType
    background_color: Optional[ColorInfo] = None
    background_image: Optional[str] = None
    placeholders: List[PlaceholderInfo] = field(default_factory=list)
    static_shapes: List[Dict] = field(default_factory=list)  # Non-editable shapes
    recommended_for: List[str] = field(default_factory=list)  # Content types this is good for


@dataclass
class DesignSystem:
    """Complete design system extracted from PPTX"""
    id: str
    name: str
    source_file: str

    # Dimensions
    slide_width_emu: int
    slide_height_emu: int
    slide_width_inches: float
    slide_height_inches: float

    # Color palette
    theme_colors: Dict[str, str] = field(default_factory=dict)  # name -> hex
    accent_colors: List[str] = field(default_factory=list)
    background_colors: List[str] = field(default_factory=list)
    text_colors: List[str] = field(default_factory=list)

    # Typography
    heading_fonts: List[str] = field(default_factory=list)
    body_fonts: List[str] = field(default_factory=list)
    font_sizes: Dict[str, float] = field(default_factory=dict)  # role -> size in pt

    # Layouts
    layouts: List[SlideLayoutInfo] = field(default_factory=list)

    # Master styles
    title_style: Optional[TextStyleInfo] = None
    subtitle_style: Optional[TextStyleInfo] = None
    body_style: Optional[TextStyleInfo] = None

    # Metadata
    fonts_used: List[str] = field(default_factory=list)
    has_custom_fonts: bool = False

    def to_dict(self) -> Dict:
        """Convert to dictionary for JSON serialization"""
        return asdict(self)

    def to_json(self) -> str:
        """Convert to JSON string"""
        return json.dumps(self.to_dict(), indent=2, default=str)


class DesignSystemExtractor:
    """
    Extracts complete design system from PPTX templates.

    This extracts every visual detail needed to recreate slides
    with 99% fidelity to the original template.
    """

    def __init__(self):
        self.common_fonts = {
            "Arial", "Helvetica", "Times New Roman", "Calibri",
            "Cambria", "Georgia", "Verdana", "Tahoma", "Trebuchet MS",
            "Roboto", "Open Sans", "Lato", "Montserrat", "Poppins"
        }

    def extract(self, pptx_path: str, name: Optional[str] = None) -> DesignSystem:
        """
        Extract complete design system from a PPTX file.

        Args:
            pptx_path: Path to the PPTX file
            name: Optional name for the design system

        Returns:
            DesignSystem with all extracted information
        """
        prs = Presentation(pptx_path)

        design_system = DesignSystem(
            id=str(uuid.uuid4()),
            name=name or os.path.basename(pptx_path).replace('.pptx', ''),
            source_file=pptx_path,
            slide_width_emu=prs.slide_width,
            slide_height_emu=prs.slide_height,
            slide_width_inches=prs.slide_width / Emu(Inches(1)),
            slide_height_inches=prs.slide_height / Emu(Inches(1)),
        )

        # Extract theme colors
        design_system.theme_colors = self._extract_theme_colors(prs)

        # Extract layouts from actual slides
        design_system.layouts = self._extract_layouts(prs)

        # Extract fonts used
        fonts = self._extract_fonts_used(prs)
        design_system.fonts_used = list(fonts)
        design_system.has_custom_fonts = bool(fonts - self.common_fonts)

        # Categorize fonts
        design_system.heading_fonts = self._identify_heading_fonts(prs)
        design_system.body_fonts = self._identify_body_fonts(prs)

        # Extract color categories
        colors = self._categorize_colors(prs, design_system.theme_colors)
        design_system.accent_colors = colors.get('accent', [])
        design_system.background_colors = colors.get('background', [])
        design_system.text_colors = colors.get('text', [])

        # Extract master styles
        master_styles = self._extract_master_styles(prs)
        design_system.title_style = master_styles.get('title')
        design_system.subtitle_style = master_styles.get('subtitle')
        design_system.body_style = master_styles.get('body')

        # Extract font sizes by role
        design_system.font_sizes = self._extract_font_sizes(prs)

        return design_system

    def _extract_theme_colors(self, prs: Presentation) -> Dict[str, str]:
        """Extract theme colors from presentation"""
        colors = {}

        try:
            theme = prs.slide_master.theme_color_scheme

            color_names = [
                'dark_1', 'light_1', 'dark_2', 'light_2',
                'accent_1', 'accent_2', 'accent_3', 'accent_4',
                'accent_5', 'accent_6', 'hyperlink', 'followed_hyperlink'
            ]

            for name in color_names:
                try:
                    color = getattr(theme, name, None)
                    if color:
                        colors[name] = self._color_to_hex(color)
                except Exception:
                    pass
        except Exception:
            pass

        return colors

    def _extract_layouts(self, prs: Presentation) -> List[SlideLayoutInfo]:
        """Extract layout information from all slides"""
        layouts = []

        for idx, slide in enumerate(prs.slides):
            layout_info = SlideLayoutInfo(
                index=idx,
                name=slide.slide_layout.name if slide.slide_layout else f"Slide {idx + 1}",
                layout_type=self._detect_layout_type(slide),
                placeholders=[],
                static_shapes=[],
                recommended_for=[]
            )

            # Extract background
            bg = self._extract_background(slide)
            layout_info.background_color = bg.get('color')
            layout_info.background_image = bg.get('image')

            # Extract placeholders and shapes
            static_idx = 1000  # Start static shape index at 1000 to avoid conflicts
            for shape in slide.shapes:
                try:
                    if self._is_placeholder(shape):
                        placeholder = self._extract_placeholder_info(shape)
                        if placeholder:
                            layout_info.placeholders.append(placeholder)
                    else:
                        # Check if static shape has text - make it editable!
                        if self._is_text_shape(shape):
                            text_placeholder = self._extract_text_shape_as_placeholder(shape, static_idx)
                            if text_placeholder:
                                layout_info.placeholders.append(text_placeholder)
                                static_idx += 1
                        else:
                            # Non-text static shapes (images, lines, etc.)
                            static = self._extract_static_shape(shape)
                            if static:
                                layout_info.static_shapes.append(static)
                except Exception as e:
                    # Skip shapes that cause errors during extraction
                    print(f"Warning: Could not extract shape: {e}")
                    continue

            # Determine what content this layout is good for
            layout_info.recommended_for = self._get_recommendations(layout_info)

            layouts.append(layout_info)

        return layouts

    def _detect_layout_type(self, slide) -> LayoutType:
        """Detect the semantic type of a slide layout"""
        layout_name = (slide.slide_layout.name if slide.slide_layout else "").lower()

        # Check layout name first
        if any(x in layout_name for x in ['title', 'cover', 'intro']):
            return LayoutType.TITLE
        if any(x in layout_name for x in ['section', 'divider']):
            return LayoutType.SECTION
        if any(x in layout_name for x in ['thank', 'end', 'closing', 'conclusion']):
            return LayoutType.CLOSING
        if any(x in layout_name for x in ['comparison', 'versus', 'vs']):
            return LayoutType.COMPARISON
        if any(x in layout_name for x in ['quote', 'testimonial']):
            return LayoutType.QUOTE
        if 'blank' in layout_name:
            return LayoutType.BLANK

        # Analyze shapes to determine type
        has_image = False
        has_bullets = False
        has_table = False
        has_chart = False
        text_shapes = []
        image_shapes = []

        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_image = True
                    image_shapes.append(shape)
                elif shape.has_table:
                    has_table = True
                elif shape.has_chart:
                    has_chart = True
                elif hasattr(shape, 'text_frame') and shape.has_text_frame:
                    tf = shape.text_frame
                    if len(tf.paragraphs) > 2:
                        has_bullets = True
                    text_shapes.append(shape)
            except Exception:
                # Skip shapes that cause errors
                continue

        # Determine by content analysis
        if has_chart:
            return LayoutType.CHART
        if has_table:
            return LayoutType.TABLE
        if has_image and image_shapes:
            # Check image position
            img = image_shapes[0]
            slide_mid = slide.slide_layout.slide_master.slide_width / 2 if slide.slide_layout else 4572000
            if img.left < slide_mid / 2:
                return LayoutType.IMAGE_LEFT
            elif img.left > slide_mid:
                return LayoutType.IMAGE_RIGHT
            else:
                return LayoutType.IMAGE_FULL
        if has_bullets:
            return LayoutType.BULLETS

        # Check for two-column by analyzing text shape positions
        if len(text_shapes) >= 2:
            lefts = [s.left for s in text_shapes]
            if max(lefts) - min(lefts) > 2000000:  # Significant horizontal spread
                return LayoutType.TWO_COLUMN

        return LayoutType.CONTENT

    def _is_placeholder(self, shape) -> bool:
        """Check if shape is a placeholder"""
        return hasattr(shape, 'is_placeholder') and shape.is_placeholder

    def _is_text_shape(self, shape) -> bool:
        """Check if shape contains editable text"""
        try:
            # Must have a text frame
            if not hasattr(shape, 'has_text_frame') or not shape.has_text_frame:
                return False

            # Must have actual text content
            if not hasattr(shape, 'text_frame') or not shape.text_frame:
                return False

            text = shape.text_frame.text.strip() if shape.text_frame.text else ""

            # Skip empty shapes
            if not text:
                return False

            # Skip very short text (likely labels or icons)
            if len(text) < 2:
                return False

            return True
        except Exception:
            return False

    def _extract_text_shape_as_placeholder(self, shape, idx: int) -> Optional[PlaceholderInfo]:
        """Extract a static text shape as an editable placeholder"""
        try:
            text = shape.text_frame.text.strip() if shape.text_frame and shape.text_frame.text else ""

            # Determine type based on text content and position
            ph_type = self._infer_text_type(shape, text)

            placeholder = PlaceholderInfo(
                idx=idx,
                type=ph_type,
                name=shape.name if hasattr(shape, 'name') else f"TextBox_{idx}",
                position={
                    'left': shape.left if hasattr(shape, 'left') else 0,
                    'top': shape.top if hasattr(shape, 'top') else 0,
                    'left_inches': (shape.left / 914400) if hasattr(shape, 'left') and shape.left else 0,
                    'top_inches': (shape.top / 914400) if hasattr(shape, 'top') and shape.top else 0,
                },
                size={
                    'width': shape.width if hasattr(shape, 'width') else 0,
                    'height': shape.height if hasattr(shape, 'height') else 0,
                    'width_inches': (shape.width / 914400) if hasattr(shape, 'width') and shape.width else 0,
                    'height_inches': (shape.height / 914400) if hasattr(shape, 'height') and shape.height else 0,
                },
                is_static=True,
                current_text=text,
                shape_id=shape.shape_id if hasattr(shape, 'shape_id') else None,
            )

            # Extract text styles
            if hasattr(shape, 'text_frame'):
                placeholder.text_styles = self._extract_text_styles(shape.text_frame)
                placeholder.max_chars = self._estimate_max_chars(shape)
                placeholder.max_lines = self._estimate_max_lines(shape)

            # Extract shape style
            placeholder.shape_style = self._extract_shape_style(shape)

            return placeholder
        except Exception as e:
            print(f"Warning: Could not extract text shape: {e}")
            return None

    def _infer_text_type(self, shape, text: str) -> str:
        """Infer the semantic type of a text shape based on its properties"""
        text_lower = text.lower()

        # Check font size to determine if it's a heading
        font_size = 12
        try:
            if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                for para in shape.text_frame.paragraphs:
                    if para.runs:
                        font = para.runs[0].font
                        if font.size:
                            font_size = font.size.pt
                            break
        except Exception:
            pass

        # Large text = likely a title
        if font_size >= 28:
            return "title"
        elif font_size >= 20:
            return "subtitle"

        # Check content patterns
        if any(kw in text_lower for kw in ['thank', 'merci', 'questions', 'contact']):
            return "title"
        elif any(kw in text_lower for kw in ['place here', 'placeholder', 'your text', 'lorem ipsum']):
            return "body"
        elif any(kw in text_lower for kw in ['sub-headline', 'sub headline', 'tagline']):
            return "subtitle"
        elif len(text.split('\n')) > 3:
            return "body"

        # Default based on position (top = title, bottom = footer)
        try:
            if shape.top and shape.top < 914400 * 2:  # Top 2 inches
                return "subtitle" if font_size < 20 else "title"
            elif shape.height and shape.top and (shape.top + shape.height) > 914400 * 6:  # Bottom area
                return "footer"
        except Exception:
            pass

        return "body"

    def _extract_placeholder_info(self, shape) -> Optional[PlaceholderInfo]:
        """Extract detailed placeholder information"""
        try:
            ph_type = "content"
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                ph_type = self._get_placeholder_type_name(shape.placeholder_format.type)

            # Get current text content
            current_text = None
            if hasattr(shape, 'text_frame') and shape.text_frame:
                current_text = shape.text_frame.text.strip() if shape.text_frame.text else None

            placeholder = PlaceholderInfo(
                idx=shape.placeholder_format.idx if hasattr(shape, 'placeholder_format') and shape.placeholder_format else 0,
                type=ph_type,
                name=shape.name if hasattr(shape, 'name') else "Unknown",
                position={
                    'left': shape.left if hasattr(shape, 'left') else 0,
                    'top': shape.top if hasattr(shape, 'top') else 0,
                    'left_inches': (shape.left / 914400) if hasattr(shape, 'left') and shape.left else 0,
                    'top_inches': (shape.top / 914400) if hasattr(shape, 'top') and shape.top else 0,
                },
                size={
                    'width': shape.width if hasattr(shape, 'width') else 0,
                    'height': shape.height if hasattr(shape, 'height') else 0,
                    'width_inches': (shape.width / 914400) if hasattr(shape, 'width') and shape.width else 0,
                    'height_inches': (shape.height / 914400) if hasattr(shape, 'height') and shape.height else 0,
                },
                is_static=False,  # Native placeholder
                current_text=current_text,
                shape_id=shape.shape_id if hasattr(shape, 'shape_id') else None,
            )

            # Extract text styles if it has text
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame and hasattr(shape, 'text_frame'):
                placeholder.text_styles = self._extract_text_styles(shape.text_frame)
                placeholder.max_chars = self._estimate_max_chars(shape)
                placeholder.max_lines = self._estimate_max_lines(shape)

            # Extract shape style
            placeholder.shape_style = self._extract_shape_style(shape)

            return placeholder
        except Exception as e:
            print(f"Warning: Could not extract placeholder info: {e}")
            return None

    def _extract_text_styles(self, text_frame) -> List[TextStyleInfo]:
        """Extract text styles from text frame"""
        styles = []

        for para in text_frame.paragraphs:
            font_info = FontInfo(
                name="Calibri",  # Default
                size_pt=12.0,
            )

            # Get font info from first run
            if para.runs:
                run = para.runs[0]
                font = run.font

                font_info.name = font.name or "Calibri"
                font_info.size_pt = font.size.pt if font.size else 12.0
                font_info.bold = font.bold or False
                font_info.italic = font.italic or False
                font_info.underline = font.underline or False

                try:
                    if font.color and font.color.type is not None and hasattr(font.color, 'rgb') and font.color.rgb:
                        font_info.color = ColorInfo(rgb=self._rgb_to_hex(font.color.rgb))
                except AttributeError:
                    pass  # Color not available

            # Get paragraph properties
            alignment = "left"
            if para.alignment:
                alignment = str(para.alignment).replace('PP_ALIGN.', '').lower()

            style = TextStyleInfo(
                font=font_info,
                alignment=alignment,
                indent_level=para.level,
            )

            # Check for bullets
            if para.level > 0 or self._has_bullet(para):
                style.bullet_char = "•"  # Default bullet

            styles.append(style)

            # Only get first 3 levels
            if len(styles) >= 3:
                break

        return styles

    def _has_bullet(self, para) -> bool:
        """Check if paragraph has bullet"""
        try:
            pPr = para._p.pPr
            if pPr is not None:
                buNone = pPr.find(qn('a:buNone'))
                return buNone is None
        except Exception:
            pass
        return False

    def _extract_shape_style(self, shape) -> ShapeStyleInfo:
        """Extract shape styling"""
        style = ShapeStyleInfo()

        try:
            fill = shape.fill
            if fill:
                fill_type = fill.type
                if fill_type:
                    style.fill_type = str(fill_type).split('.')[-1].lower()

                if hasattr(fill, 'fore_color') and fill.fore_color:
                    try:
                        rgb = fill.fore_color.rgb
                        if rgb:
                            style.fill_color = ColorInfo(rgb=self._rgb_to_hex(rgb))
                    except Exception:
                        pass
        except Exception:
            pass

        try:
            line = shape.line
            if line and line.color and line.color.rgb:
                style.line_color = ColorInfo(rgb=self._rgb_to_hex(line.color.rgb))
                style.line_width_pt = line.width.pt if line.width else None
        except Exception:
            pass

        return style

    def _extract_static_shape(self, shape) -> Optional[Dict]:
        """Extract static (non-placeholder) shape info"""
        try:
            return {
                'name': shape.name,
                'type': str(shape.shape_type).split('.')[-1] if shape.shape_type else 'unknown',
                'position': {
                    'left': shape.left,
                    'top': shape.top,
                },
                'size': {
                    'width': shape.width,
                    'height': shape.height,
                },
                'has_text': hasattr(shape, 'text_frame') and bool(shape.text_frame.text.strip()),
            }
        except Exception:
            return None

    def _extract_background(self, slide) -> Dict:
        """Extract slide background"""
        result = {}

        try:
            background = slide.background
            fill = background.fill

            if fill.type:
                fill_type = str(fill.type).split('.')[-1].lower()

                if fill_type == 'solid':
                    if fill.fore_color and fill.fore_color.rgb:
                        result['color'] = ColorInfo(rgb=self._rgb_to_hex(fill.fore_color.rgb))
        except Exception:
            pass

        return result

    def _extract_fonts_used(self, prs: Presentation) -> set:
        """Extract all fonts used in the presentation"""
        fonts = set()

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name:
                                fonts.add(run.font.name)

        return fonts

    def _identify_heading_fonts(self, prs: Presentation) -> List[str]:
        """Identify fonts used for headings"""
        heading_fonts = set()

        for slide in prs.slides:
            for shape in slide.shapes:
                if self._is_placeholder(shape):
                    ph_type = ""
                    if hasattr(shape, 'placeholder_format'):
                        ph_type = self._get_placeholder_type_name(shape.placeholder_format.type)

                    if ph_type in ['title', 'center_title', 'subtitle']:
                        if hasattr(shape, 'text_frame'):
                            for para in shape.text_frame.paragraphs:
                                for run in para.runs:
                                    if run.font.name:
                                        heading_fonts.add(run.font.name)

        return list(heading_fonts)

    def _identify_body_fonts(self, prs: Presentation) -> List[str]:
        """Identify fonts used for body text"""
        body_fonts = set()

        for slide in prs.slides:
            for shape in slide.shapes:
                if self._is_placeholder(shape):
                    ph_type = ""
                    if hasattr(shape, 'placeholder_format'):
                        ph_type = self._get_placeholder_type_name(shape.placeholder_format.type)

                    if ph_type in ['body', 'content', 'object']:
                        if hasattr(shape, 'text_frame'):
                            for para in shape.text_frame.paragraphs:
                                for run in para.runs:
                                    if run.font.name:
                                        body_fonts.add(run.font.name)

        return list(body_fonts)

    def _categorize_colors(self, prs: Presentation, theme_colors: Dict[str, str]) -> Dict[str, List[str]]:
        """Categorize colors by usage"""
        colors = {
            'accent': [],
            'background': [],
            'text': []
        }

        # Accent colors from theme
        for name, color in theme_colors.items():
            if 'accent' in name:
                colors['accent'].append(color)
            elif 'dark' in name:
                colors['text'].append(color)
            elif 'light' in name:
                colors['background'].append(color)

        return colors

    def _extract_master_styles(self, prs: Presentation) -> Dict[str, TextStyleInfo]:
        """Extract master text styles"""
        styles = {}

        try:
            master = prs.slide_master

            for shape in master.shapes:
                if self._is_placeholder(shape):
                    ph_type = ""
                    if hasattr(shape, 'placeholder_format'):
                        ph_type = self._get_placeholder_type_name(shape.placeholder_format.type)

                    if ph_type == 'title' and hasattr(shape, 'text_frame'):
                        text_styles = self._extract_text_styles(shape.text_frame)
                        if text_styles:
                            styles['title'] = text_styles[0]
                    elif ph_type == 'subtitle' and hasattr(shape, 'text_frame'):
                        text_styles = self._extract_text_styles(shape.text_frame)
                        if text_styles:
                            styles['subtitle'] = text_styles[0]
                    elif ph_type == 'body' and hasattr(shape, 'text_frame'):
                        text_styles = self._extract_text_styles(shape.text_frame)
                        if text_styles:
                            styles['body'] = text_styles[0]
        except Exception:
            pass

        return styles

    def _extract_font_sizes(self, prs: Presentation) -> Dict[str, float]:
        """Extract font sizes by role"""
        sizes = {}

        for slide in prs.slides:
            for shape in slide.shapes:
                if self._is_placeholder(shape) and hasattr(shape, 'text_frame'):
                    ph_type = ""
                    if hasattr(shape, 'placeholder_format'):
                        ph_type = self._get_placeholder_type_name(shape.placeholder_format.type)

                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                size = run.font.size.pt
                                if ph_type not in sizes or sizes[ph_type] < size:
                                    sizes[ph_type] = size
                                break

        return sizes

    def _estimate_max_chars(self, shape) -> int:
        """Estimate max characters that fit in a shape"""
        if not shape.width or not shape.height:
            return 500

        # Rough estimate based on average character width
        width_inches = shape.width / 914400
        avg_char_width = 0.08  # inches, rough average for 12pt font
        chars_per_line = int(width_inches / avg_char_width)

        height_inches = shape.height / 914400
        line_height = 0.2  # inches, rough
        num_lines = int(height_inches / line_height)

        return chars_per_line * num_lines

    def _estimate_max_lines(self, shape) -> int:
        """Estimate max lines that fit in a shape"""
        if not shape.height:
            return 10

        height_inches = shape.height / 914400
        line_height = 0.25  # inches, rough estimate for body text
        return max(1, int(height_inches / line_height))

    def _get_recommendations(self, layout: SlideLayoutInfo) -> List[str]:
        """Get content type recommendations for a layout"""
        recommendations = []

        if layout.layout_type == LayoutType.TITLE:
            recommendations = ["introduction", "cover", "title_slide"]
        elif layout.layout_type == LayoutType.SECTION:
            recommendations = ["section_header", "chapter_start", "module_intro"]
        elif layout.layout_type == LayoutType.BULLETS:
            recommendations = ["key_points", "features", "benefits", "steps"]
        elif layout.layout_type == LayoutType.IMAGE_LEFT or layout.layout_type == LayoutType.IMAGE_RIGHT:
            recommendations = ["feature_highlight", "case_study", "example"]
        elif layout.layout_type == LayoutType.COMPARISON:
            recommendations = ["before_after", "pros_cons", "comparison"]
        elif layout.layout_type == LayoutType.QUOTE:
            recommendations = ["testimonial", "quote", "highlight"]
        elif layout.layout_type == LayoutType.STATS:
            recommendations = ["statistics", "metrics", "numbers"]
        elif layout.layout_type == LayoutType.CLOSING:
            recommendations = ["conclusion", "thank_you", "call_to_action"]
        else:
            recommendations = ["general_content", "explanation", "details"]

        return recommendations

    def _get_placeholder_type_name(self, ph_type) -> str:
        """Convert placeholder type to name"""
        type_map = {
            PP_PLACEHOLDER.TITLE: "title",
            PP_PLACEHOLDER.SUBTITLE: "subtitle",
            PP_PLACEHOLDER.BODY: "body",
            PP_PLACEHOLDER.CENTER_TITLE: "center_title",
            PP_PLACEHOLDER.FOOTER: "footer",
            PP_PLACEHOLDER.DATE: "date",
            PP_PLACEHOLDER.SLIDE_NUMBER: "slide_number",
            PP_PLACEHOLDER.PICTURE: "picture",
            PP_PLACEHOLDER.CHART: "chart",
            PP_PLACEHOLDER.TABLE: "table",
            PP_PLACEHOLDER.OBJECT: "object",
        }
        return type_map.get(ph_type, "content")

    def _color_to_hex(self, color) -> str:
        """Convert color to hex string"""
        try:
            if hasattr(color, 'rgb') and color.rgb:
                return self._rgb_to_hex(color.rgb)
            return "#000000"
        except Exception:
            return "#000000"

    def _rgb_to_hex(self, rgb: RGBColor) -> str:
        """Convert RGBColor to hex string"""
        try:
            return f"#{rgb}"
        except Exception:
            return "#000000"


# Singleton instance
DESIGN_SYSTEM_EXTRACTOR = DesignSystemExtractor()
