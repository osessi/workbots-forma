"""
PPTX Template Service - Professional PowerPoint Template Management

This service handles:
1. Template upload and storage
2. Placeholder detection and mapping
3. Content replacement while preserving original design
4. High-quality export generation
"""

import os
import uuid
import shutil
import tempfile
from typing import Optional, List, Dict, Any, Tuple
from datetime import datetime
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.shapes.base import BaseShape
from pptx.shapes.picture import Picture
from pptx.shapes.placeholder import PlaceholderPicture
from PIL import Image
import aiohttp
import asyncio

from utils.get_env import get_app_data_directory_env as get_app_data_directory


class PptxTemplateService:
    """
    Service for managing PPTX templates and generating high-quality exports.
    """

    def __init__(self):
        self.templates_dir = os.path.join(get_app_data_directory(), "pptx_templates")
        self.exports_dir = os.path.join(get_app_data_directory(), "exports")
        self._ensure_directories()

    def _ensure_directories(self):
        """Create necessary directories if they don't exist."""
        os.makedirs(self.templates_dir, exist_ok=True)
        os.makedirs(self.exports_dir, exist_ok=True)

    # ==================== TEMPLATE ANALYSIS ====================

    def analyze_template(self, pptx_path: str) -> Dict[str, Any]:
        """
        Analyze a PPTX template to extract metadata, placeholders, and structure.

        Returns detailed information about:
        - Slide count and layouts
        - Placeholder shapes and their positions
        - Fonts and colors used
        - Theme information
        """
        prs = Presentation(pptx_path)

        analysis = {
            "slide_count": len(prs.slides),
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
            "slides": [],
            "fonts": set(),
            "theme_colors": {},
            "placeholder_mapping": [],
            "slide_layouts": {}
        }

        # Analyze each slide
        for slide_idx, slide in enumerate(prs.slides):
            slide_info = {
                "index": slide_idx,
                "layout_name": slide.slide_layout.name if slide.slide_layout else "Unknown",
                "shapes": [],
                "placeholders": []
            }

            for shape in slide.shapes:
                shape_info = self._analyze_shape(shape, slide_idx)
                slide_info["shapes"].append(shape_info)

                # Collect fonts
                if shape_info.get("fonts"):
                    analysis["fonts"].update(shape_info["fonts"])

                # Track placeholders
                if shape_info.get("is_placeholder") or shape_info.get("has_text"):
                    placeholder_entry = {
                        "slide_index": slide_idx,
                        "shape_id": shape.shape_id,
                        "shape_name": shape.name,
                        "placeholder_type": shape_info.get("placeholder_type", "content"),
                        "shape_type": shape_info.get("shape_type"),
                        "position": shape_info.get("position"),
                        "size": shape_info.get("size"),
                        "default_text": shape_info.get("text", "")[:200],  # First 200 chars
                        "text_format": shape_info.get("text_format")
                    }
                    analysis["placeholder_mapping"].append(placeholder_entry)
                    slide_info["placeholders"].append(placeholder_entry)

            analysis["slides"].append(slide_info)

            # Auto-detect slide purpose based on layout name and content
            layout_name = slide_info["layout_name"].lower()
            if any(x in layout_name for x in ["title", "intro", "cover"]):
                analysis["slide_layouts"]["intro"] = slide_idx
            elif any(x in layout_name for x in ["content", "text", "body"]):
                if "content" not in analysis["slide_layouts"]:
                    analysis["slide_layouts"]["content"] = slide_idx
            elif any(x in layout_name for x in ["bullet", "list"]):
                analysis["slide_layouts"]["bullets"] = slide_idx
            elif any(x in layout_name for x in ["end", "thank", "conclusion"]):
                analysis["slide_layouts"]["conclusion"] = slide_idx
            elif any(x in layout_name for x in ["section", "divider"]):
                analysis["slide_layouts"]["section"] = slide_idx

        # Extract theme colors
        analysis["theme_colors"] = self._extract_theme_colors(prs)
        analysis["fonts"] = list(analysis["fonts"])

        return analysis

    def _analyze_shape(self, shape: BaseShape, slide_idx: int) -> Dict[str, Any]:
        """Analyze a single shape and extract its properties."""
        shape_info = {
            "shape_id": shape.shape_id,
            "name": shape.name,
            "shape_type": str(shape.shape_type) if hasattr(shape, 'shape_type') else "unknown",
            "position": {
                "left": shape.left,
                "top": shape.top
            },
            "size": {
                "width": shape.width,
                "height": shape.height
            },
            "is_placeholder": False,
            "has_text": False,
            "fonts": set()
        }

        # Check if it's a placeholder
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            shape_info["is_placeholder"] = True
            if hasattr(shape, 'placeholder_format'):
                ph_type = shape.placeholder_format.type
                shape_info["placeholder_type"] = self._get_placeholder_type_name(ph_type)

        # Extract text content and formatting
        if hasattr(shape, 'text_frame'):
            try:
                tf = shape.text_frame
                if tf.text.strip():
                    shape_info["has_text"] = True
                    shape_info["text"] = tf.text

                    # Extract text formatting from first paragraph
                    if tf.paragraphs:
                        para = tf.paragraphs[0]
                        shape_info["text_format"] = {
                            "alignment": str(para.alignment) if para.alignment else "left",
                            "level": para.level
                        }

                        # Get font info from first run
                        if para.runs:
                            run = para.runs[0]
                            font = run.font
                            if font.name:
                                shape_info["fonts"].add(font.name)
                            shape_info["text_format"]["font_name"] = font.name
                            shape_info["text_format"]["font_size"] = font.size.pt if font.size else None
                            shape_info["text_format"]["bold"] = font.bold
                            shape_info["text_format"]["italic"] = font.italic
                            if font.color and font.color.rgb:
                                shape_info["text_format"]["color"] = str(font.color.rgb)
            except Exception:
                pass

        # Check for images
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape_info["is_picture"] = True

        # Check for tables
        if hasattr(shape, 'table'):
            shape_info["is_table"] = True
            shape_info["table_size"] = {
                "rows": len(shape.table.rows),
                "cols": len(shape.table.columns)
            }

        # Check for charts
        if hasattr(shape, 'chart'):
            shape_info["is_chart"] = True
            shape_info["chart_type"] = str(shape.chart.chart_type) if shape.chart else None

        return shape_info

    def _get_placeholder_type_name(self, ph_type) -> str:
        """Convert placeholder type enum to readable name."""
        type_map = {
            PP_PLACEHOLDER.TITLE: "title",
            PP_PLACEHOLDER.SUBTITLE: "subtitle",
            PP_PLACEHOLDER.BODY: "body",
            PP_PLACEHOLDER.CENTER_TITLE: "center_title",
            PP_PLACEHOLDER.FOOTER: "footer",
            PP_PLACEHOLDER.DATE: "date",
            PP_PLACEHOLDER.SLIDE_NUMBER: "slide_number",
            PP_PLACEHOLDER.PICTURE: "picture",
            PP_PLACEHOLDER.CHART: "chart",
            PP_PLACEHOLDER.TABLE: "table",
            PP_PLACEHOLDER.OBJECT: "object",
        }
        return type_map.get(ph_type, "content")

    def _extract_theme_colors(self, prs: Presentation) -> Dict[str, str]:
        """Extract theme colors from the presentation."""
        colors = {}
        try:
            theme = prs.slide_master.theme_color_scheme
            if theme:
                for color_name in ["accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
                                   "dark1", "dark2", "light1", "light2"]:
                    try:
                        color = getattr(theme, color_name, None)
                        if color:
                            colors[color_name] = str(color)
                    except Exception:
                        pass
        except Exception:
            pass
        return colors

    # ==================== TEMPLATE UPLOAD ====================

    async def save_template(
        self,
        file_content: bytes,
        name: str,
        description: Optional[str] = None,
        category: str = "general"
    ) -> Dict[str, Any]:
        """
        Save an uploaded PPTX template and analyze it.

        Returns template metadata including placeholder mapping.
        """
        template_id = str(uuid.uuid4())

        # Create template directory
        template_dir = os.path.join(self.templates_dir, template_id)
        os.makedirs(template_dir, exist_ok=True)

        # Save the original PPTX file
        pptx_path = os.path.join(template_dir, "template.pptx")
        with open(pptx_path, "wb") as f:
            f.write(file_content)

        # Analyze the template
        analysis = self.analyze_template(pptx_path)

        # Generate thumbnail from first slide
        thumbnail_path = await self._generate_thumbnail(pptx_path, template_dir)

        return {
            "id": template_id,
            "name": name,
            "description": description,
            "category": category,
            "file_path": pptx_path,
            "file_size": len(file_content),
            "slide_count": analysis["slide_count"],
            "thumbnail_path": thumbnail_path,
            "placeholder_mapping": analysis["placeholder_mapping"],
            "slide_layouts": analysis["slide_layouts"],
            "fonts": analysis["fonts"],
            "theme_colors": analysis["theme_colors"],
            "analysis": analysis
        }

    async def _generate_thumbnail(self, pptx_path: str, output_dir: str) -> Optional[str]:
        """Generate a thumbnail image from the first slide."""
        try:
            # This would require a PPTX to image converter
            # For now, return None - can be implemented with libreoffice or similar
            return None
        except Exception:
            return None

    # ==================== CONTENT REPLACEMENT ENGINE ====================

    async def generate_from_template(
        self,
        template_path: str,
        slides_content: List[Dict[str, Any]],
        output_filename: str,
        options: Optional[Dict[str, Any]] = None
    ) -> str:
        """
        Generate a new PPTX from a template by replacing content.

        This is the core function that:
        1. Loads the original template
        2. Duplicates/removes slides as needed
        3. Replaces text in placeholders
        4. Replaces/adds images
        5. Preserves all original styling

        Args:
            template_path: Path to the template PPTX
            slides_content: List of content for each slide
            output_filename: Name for the output file
            options: Additional options (fonts, colors override, etc.)

        Returns:
            Path to the generated PPTX file
        """
        options = options or {}

        # Load the template
        prs = Presentation(template_path)
        template_slide_count = len(prs.slides)

        # Determine how to map content to slides
        content_count = len(slides_content)

        if content_count <= template_slide_count:
            # We have enough slides, just fill them
            for i, content in enumerate(slides_content):
                await self._fill_slide(prs.slides[i], content, options)

            # Remove unused slides (from end)
            for _ in range(template_slide_count - content_count):
                self._remove_slide(prs, len(prs.slides) - 1)
        else:
            # Need to duplicate slides
            # First, fill all existing slides
            for i in range(template_slide_count):
                if i < content_count:
                    await self._fill_slide(prs.slides[i], slides_content[i], options)

            # Then duplicate appropriate slides for remaining content
            content_slide_idx = 1 if template_slide_count > 1 else 0  # Use second slide as template
            for i in range(template_slide_count, content_count):
                new_slide = self._duplicate_slide(prs, content_slide_idx)
                await self._fill_slide(new_slide, slides_content[i], options)

        # Save the output
        output_path = os.path.join(self.exports_dir, f"{output_filename}.pptx")
        prs.save(output_path)

        return output_path

    async def _fill_slide(
        self,
        slide,
        content: Dict[str, Any],
        options: Dict[str, Any]
    ):
        """
        Fill a slide with content while preserving original formatting.

        Content structure:
        {
            "title": "Slide Title",
            "subtitle": "Optional subtitle",
            "body": "Main content text",
            "bullets": ["Point 1", "Point 2", ...],
            "image": {"url": "...", "prompt": "..."},
            "table": [[...], [...], ...],
            "speaker_notes": "Notes for presenter"
        }
        """
        # Build a mapping of placeholder types to shapes
        shape_map = {}
        for shape in slide.shapes:
            if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                ph_type = self._get_placeholder_type_name(shape.placeholder_format.type)
                shape_map[ph_type] = shape
            elif hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                # Non-placeholder text shapes - map by position
                name = shape.name.lower()
                if "title" in name:
                    shape_map.setdefault("title", shape)
                elif "subtitle" in name or "sub" in name:
                    shape_map.setdefault("subtitle", shape)
                elif "body" in name or "content" in name or "text" in name:
                    shape_map.setdefault("body", shape)

        # Also check for the first large text shape as title if not found
        if "title" not in shape_map:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.width > Inches(5):
                    shape_map["title"] = shape
                    break

        # Replace title
        if "title" in content and "title" in shape_map:
            self._replace_text_preserve_format(shape_map["title"], content["title"])
        elif "title" in content:
            # Find any text shape that might be a title
            self._find_and_replace_title(slide, content["title"])

        # Replace subtitle
        if "subtitle" in content and "subtitle" in shape_map:
            self._replace_text_preserve_format(shape_map["subtitle"], content["subtitle"])

        # Replace body content
        if "body" in content and "body" in shape_map:
            self._replace_text_preserve_format(shape_map["body"], content["body"])

        # Replace bullets
        if "bullets" in content:
            self._replace_bullets(slide, content["bullets"], shape_map.get("body"))

        # Replace/add images
        if "image" in content:
            await self._replace_image(slide, content["image"], options)

        # Add speaker notes
        if "speaker_notes" in content and content["speaker_notes"]:
            try:
                notes_slide = slide.notes_slide
                notes_tf = notes_slide.notes_text_frame
                notes_tf.text = content["speaker_notes"]
            except Exception:
                pass

    def _replace_text_preserve_format(self, shape, new_text: str):
        """
        Replace text in a shape while preserving original formatting.

        This is crucial for maintaining the template's visual quality.
        """
        if not hasattr(shape, 'text_frame'):
            return

        tf = shape.text_frame

        if not tf.paragraphs:
            return

        # For simple single-line replacement
        if "\n" not in new_text and len(tf.paragraphs) == 1:
            para = tf.paragraphs[0]
            if para.runs:
                # Keep first run's formatting, replace text
                para.runs[0].text = new_text
                # Remove other runs
                for run in list(para.runs)[1:]:
                    run.text = ""
            else:
                para.text = new_text
            return

        # For multi-paragraph text
        lines = new_text.split("\n")

        # Get formatting from first paragraph/run as template
        template_para = tf.paragraphs[0]
        template_font = None
        if template_para.runs:
            template_run = template_para.runs[0]
            template_font = {
                "name": template_run.font.name,
                "size": template_run.font.size,
                "bold": template_run.font.bold,
                "italic": template_run.font.italic,
                "color": template_run.font.color.rgb if template_run.font.color and template_run.font.color.rgb else None
            }

        # Clear existing paragraphs (except first) and fill with new content
        # Note: We can't delete paragraphs in python-pptx, so we clear them
        for i, para in enumerate(tf.paragraphs):
            if i < len(lines):
                if para.runs:
                    para.runs[0].text = lines[i]
                    for run in list(para.runs)[1:]:
                        run.text = ""
                else:
                    para.text = lines[i]
            else:
                # Clear extra paragraphs
                if para.runs:
                    for run in para.runs:
                        run.text = ""
                else:
                    para.text = ""

        # Add new paragraphs if needed
        for i in range(len(tf.paragraphs), len(lines)):
            para = tf.paragraphs[0]._p.addnext(tf.paragraphs[0]._p.makeelement(qn('a:p')))
            # This is complex in python-pptx, simplified version
            pass

    def _find_and_replace_title(self, slide, new_title: str):
        """Find the most likely title shape and replace its text."""
        # Look for shapes with "title" in name or largest text at top
        candidates = []

        for shape in slide.shapes:
            if not hasattr(shape, 'text_frame'):
                continue

            score = 0
            if "title" in shape.name.lower():
                score += 100
            if shape.top < Inches(2):  # Near top
                score += 50
            if shape.width > Inches(6):  # Wide shape
                score += 30
            if shape.text_frame.text.strip():
                score += 10

            if score > 0:
                candidates.append((score, shape))

        if candidates:
            candidates.sort(reverse=True, key=lambda x: x[0])
            self._replace_text_preserve_format(candidates[0][1], new_title)

    def _replace_bullets(self, slide, bullets: List[str], body_shape=None):
        """Replace bullet points in a slide."""
        if body_shape and hasattr(body_shape, 'text_frame'):
            tf = body_shape.text_frame

            # Get template formatting from first paragraph
            template_level = 0
            template_font = None
            if tf.paragraphs:
                template_para = tf.paragraphs[0]
                template_level = template_para.level
                if template_para.runs:
                    run = template_para.runs[0]
                    template_font = run.font

            # Replace bullet content
            for i, para in enumerate(tf.paragraphs):
                if i < len(bullets):
                    bullet_text = bullets[i]
                    if isinstance(bullet_text, dict):
                        bullet_text = bullet_text.get("text", bullet_text.get("title", str(bullet_text)))

                    if para.runs:
                        para.runs[0].text = str(bullet_text)
                        for run in list(para.runs)[1:]:
                            run.text = ""
                    else:
                        para.text = str(bullet_text)
                else:
                    # Clear extra bullets
                    if para.runs:
                        for run in para.runs:
                            run.text = ""
            return

        # If no body shape found, look for bullet shapes
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                tf = shape.text_frame
                if len(tf.paragraphs) >= 2:  # Likely a bullet list
                    for i, para in enumerate(tf.paragraphs):
                        if i < len(bullets):
                            bullet_text = bullets[i]
                            if isinstance(bullet_text, dict):
                                bullet_text = bullet_text.get("text", bullet_text.get("title", str(bullet_text)))
                            if para.runs:
                                para.runs[0].text = str(bullet_text)
                    return

    async def _replace_image(self, slide, image_data: Dict[str, Any], options: Dict[str, Any]):
        """Replace or add an image to the slide."""
        image_url = image_data.get("url") or image_data.get("__image_url__")
        if not image_url:
            return

        # Download the image
        image_bytes = await self._download_image(image_url)
        if not image_bytes:
            return

        # Find picture placeholder or existing image
        picture_shape = None
        picture_position = None

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_shape = shape
                picture_position = {
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height
                }
                break
            elif hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                if shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE:
                    picture_shape = shape
                    picture_position = {
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    }
                    break

        if picture_shape and picture_position:
            # Remove old picture and add new one at same position
            sp = picture_shape._element
            sp.getparent().remove(sp)

            # Add new image
            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                tmp.write(image_bytes)
                tmp_path = tmp.name

            try:
                slide.shapes.add_picture(
                    tmp_path,
                    picture_position["left"],
                    picture_position["top"],
                    picture_position["width"],
                    picture_position["height"]
                )
            finally:
                os.unlink(tmp_path)

    async def _download_image(self, url: str) -> Optional[bytes]:
        """Download an image from URL."""
        try:
            async with aiohttp.ClientSession() as session:
                async with session.get(url, timeout=30) as response:
                    if response.status == 200:
                        return await response.read()
        except Exception as e:
            print(f"Error downloading image: {e}")
        return None

    def _duplicate_slide(self, prs: Presentation, slide_index: int):
        """Duplicate a slide in the presentation."""
        template_slide = prs.slides[slide_index]
        slide_layout = template_slide.slide_layout

        # Add new slide with same layout
        new_slide = prs.slides.add_slide(slide_layout)

        # Copy shapes from template
        for shape in template_slide.shapes:
            self._copy_shape(shape, new_slide)

        return new_slide

    def _copy_shape(self, shape, target_slide):
        """Copy a shape to another slide."""
        # This is complex in python-pptx
        # For now, we rely on the slide layout having the same structure
        pass

    def _remove_slide(self, prs: Presentation, slide_index: int):
        """Remove a slide from the presentation."""
        slide_id = prs.slides._sldIdLst[slide_index].rId
        prs.part.drop_rel(slide_id)
        del prs.slides._sldIdLst[slide_index]


# Singleton instance
PPTX_TEMPLATE_SERVICE = PptxTemplateService()
