"""
Slide Thumbnail Generator - Generate thumbnail images from PPTX slides

Uses LibreOffice in headless mode or python-pptx with Pillow to generate
thumbnail images of slides for preview purposes.
"""

import os
import subprocess
import tempfile
import shutil
from typing import List, Optional
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Emu
from PIL import Image
import io
import base64


class SlideThumbnailGenerator:
    """
    Generate thumbnail images from PPTX slides.

    Supports multiple methods:
    1. LibreOffice conversion (best quality)
    2. PDF intermediate conversion
    3. Basic placeholder-based preview (fallback)
    """

    def __init__(self, output_dir: Optional[str] = None):
        self.output_dir = output_dir or tempfile.mkdtemp(prefix="slide_thumbs_")
        os.makedirs(self.output_dir, exist_ok=True)

        # Check if LibreOffice is available
        self.libreoffice_path = self._find_libreoffice()

    def _find_libreoffice(self) -> Optional[str]:
        """Find LibreOffice executable"""
        possible_paths = [
            "/usr/local/bin/soffice",  # macOS Homebrew
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",  # macOS native
            "/usr/bin/libreoffice",  # Linux
            "/usr/bin/soffice",  # Linux alt
            "C:\\Program Files\\LibreOffice\\program\\soffice.exe",  # Windows
            "libreoffice",  # PATH
            "soffice",  # PATH
        ]

        for path in possible_paths:
            if os.path.exists(path):
                return path
            # Check if in PATH
            try:
                result = subprocess.run(
                    ["which", path] if os.name != "nt" else ["where", path],
                    capture_output=True,
                    text=True
                )
                if result.returncode == 0:
                    return result.stdout.strip()
            except Exception:
                pass

        return None

    def generate_thumbnails(
        self,
        pptx_path: str,
        template_id: str,
        width: int = 1920,
        height: int = 1080
    ) -> List[str]:
        """
        Generate thumbnail images for all slides in a PPTX.

        Args:
            pptx_path: Path to PPTX file
            template_id: Template ID for organizing thumbnails
            width: Thumbnail width in pixels
            height: Thumbnail height in pixels

        Returns:
            List of paths to generated thumbnail images
        """
        # Create template-specific directory
        template_dir = os.path.join(self.output_dir, template_id)
        os.makedirs(template_dir, exist_ok=True)

        thumbnails = []

        # Try LibreOffice method first (best quality)
        if self.libreoffice_path:
            thumbnails = self._generate_with_libreoffice(
                pptx_path, template_dir, width, height
            )

        # Fallback to basic preview if LibreOffice not available
        if not thumbnails:
            thumbnails = self._generate_basic_preview(
                pptx_path, template_dir, width, height
            )

        return thumbnails

    def _generate_with_libreoffice(
        self,
        pptx_path: str,
        output_dir: str,
        width: int,
        height: int
    ) -> List[str]:
        """Generate thumbnails using LibreOffice"""
        thumbnails = []

        try:
            # Convert PPTX to PDF first
            with tempfile.TemporaryDirectory() as temp_dir:
                # Run LibreOffice to convert to PDF
                cmd = [
                    self.libreoffice_path,
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", temp_dir,
                    pptx_path
                ]

                result = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=120
                )

                if result.returncode != 0:
                    print(f"LibreOffice conversion failed: {result.stderr}")
                    return []

                # Find the generated PDF
                pdf_name = os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
                pdf_path = os.path.join(temp_dir, pdf_name)

                if not os.path.exists(pdf_path):
                    print(f"PDF not found: {pdf_path}")
                    return []

                # Convert PDF pages to images using pdf2image or poppler
                try:
                    from pdf2image import convert_from_path

                    images = convert_from_path(
                        pdf_path,
                        dpi=72,
                        size=(width, height)
                    )

                    for idx, image in enumerate(images):
                        thumb_path = os.path.join(output_dir, f"slide_{idx}.png")
                        image.save(thumb_path, "PNG")
                        thumbnails.append(thumb_path)

                except ImportError:
                    print("pdf2image not installed, falling back to basic preview")
                    return []

        except subprocess.TimeoutExpired:
            print("LibreOffice conversion timed out")
            return []
        except Exception as e:
            print(f"LibreOffice conversion error: {e}")
            return []

        return thumbnails

    def _generate_basic_preview(
        self,
        pptx_path: str,
        output_dir: str,
        width: int,
        height: int
    ) -> List[str]:
        """Generate basic preview images showing placeholders"""
        thumbnails = []

        try:
            prs = Presentation(pptx_path)

            # Get slide dimensions
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            for idx, slide in enumerate(prs.slides):
                # Create a basic preview image
                img = Image.new('RGB', (width, height), color='white')

                # Draw a border
                from PIL import ImageDraw
                draw = ImageDraw.Draw(img)
                draw.rectangle([(0, 0), (width-1, height-1)], outline='#e0e0e0', width=1)

                # Draw placeholders as colored rectangles
                for shape in slide.shapes:
                    try:
                        # Calculate scaled position
                        left = int((shape.left / slide_width) * width)
                        top = int((shape.top / slide_height) * height)
                        right = int(((shape.left + shape.width) / slide_width) * width)
                        bottom = int(((shape.top + shape.height) / slide_height) * height)

                        # Get color based on shape type
                        color = self._get_shape_color(shape)

                        # Draw rectangle
                        draw.rectangle(
                            [(left, top), (right, bottom)],
                            fill=color,
                            outline=self._darken_color(color)
                        )

                        # Try to draw text if available
                        if hasattr(shape, 'text') and shape.text:
                            text = shape.text[:20] + "..." if len(shape.text) > 20 else shape.text
                            try:
                                # Simple text placement
                                text_x = left + 5
                                text_y = top + 5
                                if right - left > 30 and bottom - top > 15:
                                    draw.text((text_x, text_y), text, fill='#333333')
                            except Exception:
                                pass

                    except Exception as e:
                        continue

                # Save thumbnail
                thumb_path = os.path.join(output_dir, f"slide_{idx}.png")
                img.save(thumb_path, "PNG")
                thumbnails.append(thumb_path)

        except Exception as e:
            print(f"Basic preview generation error: {e}")
            return []

        return thumbnails

    def _get_shape_color(self, shape) -> str:
        """Get color for shape based on its type"""
        from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

        # Try to get fill color
        try:
            if hasattr(shape, 'fill') and shape.fill.type is not None:
                if hasattr(shape.fill, 'fore_color') and shape.fill.fore_color.rgb:
                    rgb = shape.fill.fore_color.rgb
                    return f"#{rgb}"
        except Exception:
            pass

        # Check placeholder type
        if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
            try:
                ph_type = shape.placeholder_format.type
                type_colors = {
                    PP_PLACEHOLDER.TITLE: "#e8e0ff",  # Light purple
                    PP_PLACEHOLDER.SUBTITLE: "#e0f0ff",  # Light blue
                    PP_PLACEHOLDER.BODY: "#e0ffe0",  # Light green
                    PP_PLACEHOLDER.PICTURE: "#ffe0f0",  # Light pink
                    PP_PLACEHOLDER.CHART: "#e0ffff",  # Light cyan
                    PP_PLACEHOLDER.TABLE: "#fff0e0",  # Light orange
                }
                return type_colors.get(ph_type, "#f0f0f0")
            except Exception:
                pass

        # Check shape type
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return "#ffe0f0"
            elif shape.has_table:
                return "#fff0e0"
            elif shape.has_chart:
                return "#e0ffff"
        except Exception:
            pass

        return "#f5f5f5"

    def _darken_color(self, hex_color: str) -> str:
        """Darken a hex color for border"""
        try:
            hex_color = hex_color.lstrip('#')
            r = max(0, int(hex_color[0:2], 16) - 40)
            g = max(0, int(hex_color[2:4], 16) - 40)
            b = max(0, int(hex_color[4:6], 16) - 40)
            return f"#{r:02x}{g:02x}{b:02x}"
        except Exception:
            return "#cccccc"

    def get_thumbnail_path(self, template_id: str, slide_index: int) -> Optional[str]:
        """Get path to a specific slide thumbnail"""
        thumb_path = os.path.join(
            self.output_dir,
            template_id,
            f"slide_{slide_index}.png"
        )
        return thumb_path if os.path.exists(thumb_path) else None

    def get_thumbnail_base64(self, template_id: str, slide_index: int) -> Optional[str]:
        """Get base64 encoded thumbnail for embedding in JSON"""
        thumb_path = self.get_thumbnail_path(template_id, slide_index)
        if not thumb_path:
            return None

        try:
            with open(thumb_path, "rb") as f:
                return base64.b64encode(f.read()).decode('utf-8')
        except Exception:
            return None

    def cleanup(self, template_id: Optional[str] = None):
        """Clean up generated thumbnails"""
        if template_id:
            template_dir = os.path.join(self.output_dir, template_id)
            if os.path.exists(template_dir):
                shutil.rmtree(template_dir)
        else:
            if os.path.exists(self.output_dir):
                shutil.rmtree(self.output_dir)


# Global instance
THUMBNAIL_GENERATOR = SlideThumbnailGenerator(
    output_dir=os.path.join(
        os.getenv("APP_DATA_DIRECTORY", "/tmp/slides-api-data"),
        "thumbnails"
    )
)
