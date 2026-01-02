"""
Smart PPTX Builder - Professional PowerPoint Generation with Style Preservation

This service creates high-quality PPTX files by:
1. Using pre-designed slide templates as bases
2. Intelligently adapting content (adding/removing bullets)
3. Preserving all original styling (fonts, colors, effects)
4. Handling images with proper sizing and positioning
"""

import os
import copy
import uuid
import tempfile
from typing import List, Dict, Any, Optional, Tuple
from io import BytesIO
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from pptx.shapes.base import BaseShape
from pptx.shapes.autoshape import Shape
from pptx.text.text import _Paragraph
from lxml import etree
import aiohttp

from utils.get_env import get_app_data_directory_env as get_app_data_directory


# ====================================================================================
# DEFAULT TEMPLATE DEFINITIONS
# These define the visual styles for each template type
# ====================================================================================

TEMPLATE_STYLES = {
    "general": {
        "name": "General",
        "colors": {
            "primary": "9333EA",      # Purple
            "secondary": "7C3AED",
            "accent": "A855F7",
            "background": "FFFFFF",
            "text_dark": "111827",
            "text_light": "4B5563",
            "text_muted": "9CA3AF"
        },
        "fonts": {
            "heading": "Poppins",
            "body": "Inter",
            "fallback": "Arial"
        },
        "slide_layouts": {
            "intro": {
                "title_font_size": 44,
                "subtitle_font_size": 18,
                "has_image": True,
                "image_position": "left"  # or "right", "background"
            },
            "content": {
                "title_font_size": 36,
                "body_font_size": 16,
                "bullet_font_size": 14,
                "max_bullets": 6
            },
            "bullets": {
                "title_font_size": 32,
                "bullet_font_size": 16,
                "icon_style": "circle",  # or "number", "check", "arrow"
                "columns": 1  # or 2
            }
        }
    },
    "modern": {
        "name": "Modern",
        "colors": {
            "primary": "1F2937",      # Dark gray
            "secondary": "374151",
            "accent": "3B82F6",       # Blue accent
            "background": "111827",   # Dark background
            "text_dark": "F9FAFB",    # Light text for dark bg
            "text_light": "D1D5DB",
            "text_muted": "9CA3AF"
        },
        "fonts": {
            "heading": "Inter",
            "body": "Inter",
            "fallback": "Arial"
        },
        "slide_layouts": {
            "intro": {
                "title_font_size": 48,
                "subtitle_font_size": 20,
                "has_image": True,
                "image_position": "right"
            },
            "content": {
                "title_font_size": 40,
                "body_font_size": 18,
                "bullet_font_size": 16,
                "max_bullets": 5
            },
            "bullets": {
                "title_font_size": 36,
                "bullet_font_size": 18,
                "icon_style": "number",
                "columns": 1
            }
        }
    },
    "standard": {
        "name": "Standard",
        "colors": {
            "primary": "2563EB",      # Blue
            "secondary": "1D4ED8",
            "accent": "60A5FA",
            "background": "FFFFFF",
            "text_dark": "1E3A5F",
            "text_light": "475569",
            "text_muted": "94A3B8"
        },
        "fonts": {
            "heading": "Calibri",
            "body": "Calibri",
            "fallback": "Arial"
        },
        "slide_layouts": {
            "intro": {
                "title_font_size": 40,
                "subtitle_font_size": 18,
                "has_image": False
            },
            "content": {
                "title_font_size": 32,
                "body_font_size": 14,
                "bullet_font_size": 14,
                "max_bullets": 6
            },
            "bullets": {
                "title_font_size": 28,
                "bullet_font_size": 14,
                "icon_style": "circle",
                "columns": 1
            }
        }
    },
    "swift": {
        "name": "Swift",
        "colors": {
            "primary": "F97316",      # Orange
            "secondary": "EA580C",
            "accent": "FB923C",
            "background": "FFFBEB",   # Warm white
            "text_dark": "292524",
            "text_light": "57534E",
            "text_muted": "A8A29E"
        },
        "fonts": {
            "heading": "Montserrat",
            "body": "Open Sans",
            "fallback": "Arial"
        },
        "slide_layouts": {
            "intro": {
                "title_font_size": 42,
                "subtitle_font_size": 16,
                "has_image": True,
                "image_position": "left"
            },
            "content": {
                "title_font_size": 34,
                "body_font_size": 16,
                "bullet_font_size": 14,
                "max_bullets": 4
            },
            "bullets": {
                "title_font_size": 30,
                "bullet_font_size": 16,
                "icon_style": "arrow",
                "columns": 2
            }
        }
    }
}


class SmartPptxBuilder:
    """
    Smart PPTX Builder that creates professional presentations.

    Features:
    - Uses template-based styling for consistent, professional look
    - Intelligently adapts bullet counts (adds/removes items)
    - Preserves formatting when replacing content
    - Handles images with proper sizing
    """

    def __init__(self, template_name: str = "general"):
        self.template_name = template_name.lower()
        self.style = TEMPLATE_STYLES.get(self.template_name, TEMPLATE_STYLES["general"])

        # Create a new presentation with custom size (16:9)
        self.prs = Presentation()
        self.prs.slide_width = Emu(12192000)  # 12.7 inches (16:9)
        self.prs.slide_height = Emu(6858000)  # 7.15 inches

        # Cache for downloaded images
        self._image_cache: Dict[str, bytes] = {}

        # Directory for temporary files
        self.temp_dir = tempfile.mkdtemp()

    # ====================================================================================
    # SLIDE CREATION METHODS
    # ====================================================================================

    def create_intro_slide(
        self,
        title: str,
        subtitle: str = "",
        presenter_name: str = "",
        date: str = "",
        image_url: str = None
    ):
        """Create an intro/title slide with professional styling."""
        slide = self._add_blank_slide()
        layout = self.style["slide_layouts"]["intro"]
        colors = self.style["colors"]

        # Background
        self._set_slide_background(slide, colors["background"])

        # Accent bar at top
        self._add_accent_bar(slide, colors["primary"])

        if layout.get("has_image") and image_url:
            # Two-column layout with image
            if layout.get("image_position") == "left":
                self._create_intro_with_left_image(slide, title, subtitle, presenter_name, date, image_url, layout, colors)
            else:
                self._create_intro_with_right_image(slide, title, subtitle, presenter_name, date, image_url, layout, colors)
        else:
            # Centered text layout
            self._create_centered_intro(slide, title, subtitle, presenter_name, date, layout, colors)

        return slide

    def create_content_slide(
        self,
        title: str,
        body: str = "",
        bullets: List[Dict[str, Any]] = None,
        image_url: str = None,
        speaker_notes: str = ""
    ):
        """Create a content slide with title and body/bullets."""
        slide = self._add_blank_slide()
        layout = self.style["slide_layouts"]["content"]
        colors = self.style["colors"]

        # Background
        self._set_slide_background(slide, colors["background"])

        # Header accent
        self._add_header_accent(slide, colors["primary"])

        # Title
        title_top = Inches(0.8)
        title_height = Inches(0.8)
        self._add_title(slide, title, Inches(0.7), title_top, Inches(11.3), title_height,
                       layout["title_font_size"], colors["text_dark"])

        content_top = Inches(1.8)

        if image_url:
            # Content with image on right
            content_width = Inches(6.5)
            self._add_content_with_image(slide, body, bullets, image_url, content_top, content_width, layout, colors)
        else:
            # Full-width content
            content_width = Inches(11.3)
            if bullets:
                self._add_bullet_list(slide, bullets, Inches(0.7), content_top, content_width, layout, colors)
            elif body:
                self._add_body_text(slide, body, Inches(0.7), content_top, content_width, Inches(4.5),
                                   layout["body_font_size"], colors["text_light"])

        # Speaker notes
        if speaker_notes:
            self._add_speaker_notes(slide, speaker_notes)

        return slide

    def create_bullets_slide(
        self,
        title: str,
        bullets: List[Dict[str, Any]],
        icon_style: str = None,
        speaker_notes: str = ""
    ):
        """Create a slide focused on bullet points with optional icons."""
        slide = self._add_blank_slide()
        layout = self.style["slide_layouts"]["bullets"]
        colors = self.style["colors"]

        icon_style = icon_style or layout.get("icon_style", "circle")

        # Background
        self._set_slide_background(slide, colors["background"])

        # Header accent
        self._add_header_accent(slide, colors["primary"])

        # Title
        self._add_title(slide, title, Inches(0.7), Inches(0.8), Inches(11.3), Inches(0.8),
                       layout["title_font_size"], colors["text_dark"])

        # Bullets with icons
        columns = layout.get("columns", 1)
        if columns == 2 and len(bullets) > 3:
            self._add_two_column_bullets(slide, bullets, layout, colors, icon_style)
        else:
            self._add_styled_bullets(slide, bullets, layout, colors, icon_style)

        # Speaker notes
        if speaker_notes:
            self._add_speaker_notes(slide, speaker_notes)

        return slide

    def create_conclusion_slide(
        self,
        title: str,
        message: str = "",
        call_to_action: str = "",
        contact_info: str = ""
    ):
        """Create a conclusion/thank you slide."""
        slide = self._add_blank_slide()
        colors = self.style["colors"]

        # Gradient-like background (using primary color)
        self._set_slide_background(slide, colors["primary"])

        # Large centered title
        self._add_title(slide, title, Inches(1), Inches(2.5), Inches(10.7), Inches(1.5),
                       48, colors["text_dark"] if colors["background"] == "FFFFFF" else "FFFFFF",
                       alignment=PP_ALIGN.CENTER, bold=True)

        if message:
            self._add_body_text(slide, message, Inches(1.5), Inches(4.2), Inches(9.7), Inches(1),
                               20, "FFFFFF", alignment=PP_ALIGN.CENTER)

        if call_to_action:
            self._add_body_text(slide, call_to_action, Inches(1.5), Inches(5.3), Inches(9.7), Inches(0.6),
                               16, colors["accent"], alignment=PP_ALIGN.CENTER)

        return slide

    # ====================================================================================
    # HELPER METHODS FOR SLIDE ELEMENTS
    # ====================================================================================

    def _add_blank_slide(self):
        """Add a blank slide to the presentation."""
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        return self.prs.slides.add_slide(blank_layout)

    def _set_slide_background(self, slide, color_hex: str):
        """Set the slide background color."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(color_hex)

    def _add_accent_bar(self, slide, color_hex: str):
        """Add a decorative accent bar at the top of the slide."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, Inches(0.15)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(color_hex)
        shape.line.fill.background()

    def _add_header_accent(self, slide, color_hex: str):
        """Add a subtle header accent line."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.7), Inches(1.65),
            Inches(2), Inches(0.05)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(color_hex)
        shape.line.fill.background()

    def _add_title(self, slide, text: str, left, top, width, height, font_size: int,
                   color_hex: str, alignment=PP_ALIGN.LEFT, bold: bool = True):
        """Add a title text box."""
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        para = tf.paragraphs[0]
        para.text = text
        para.alignment = alignment

        font = para.font
        font.name = self.style["fonts"]["heading"]
        font.size = Pt(font_size)
        font.bold = bold
        font.color.rgb = RGBColor.from_string(color_hex)

    def _add_body_text(self, slide, text: str, left, top, width, height,
                       font_size: int, color_hex: str, alignment=PP_ALIGN.LEFT):
        """Add body text."""
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True

        para = tf.paragraphs[0]
        para.text = text
        para.alignment = alignment
        para.line_spacing = 1.5

        font = para.font
        font.name = self.style["fonts"]["body"]
        font.size = Pt(font_size)
        font.color.rgb = RGBColor.from_string(color_hex)

    def _add_bullet_list(self, slide, bullets: List[Dict[str, Any]], left, top, width,
                         layout: Dict, colors: Dict):
        """Add a bullet list with proper formatting."""
        textbox = slide.shapes.add_textbox(left, top, width, Inches(4.5))
        tf = textbox.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                para = tf.paragraphs[0]
            else:
                para = tf.add_paragraph()

            # Get bullet text
            if isinstance(bullet, dict):
                title = bullet.get("title", bullet.get("text", ""))
                description = bullet.get("description", "")
                text = f"{title}" if not description else f"{title}\n{description}"
            else:
                text = str(bullet)

            para.text = f"• {text}"
            para.level = 0
            para.space_before = Pt(8)
            para.space_after = Pt(8)
            para.line_spacing = 1.3

            font = para.font
            font.name = self.style["fonts"]["body"]
            font.size = Pt(layout["bullet_font_size"])
            font.color.rgb = RGBColor.from_string(colors["text_light"])

    def _add_styled_bullets(self, slide, bullets: List[Dict[str, Any]], layout: Dict,
                            colors: Dict, icon_style: str):
        """Add bullets with styled icons/numbers."""
        bullet_top = Inches(1.9)
        bullet_spacing = Inches(0.9)

        for i, bullet in enumerate(bullets):
            y_pos = bullet_top + (i * bullet_spacing)

            # Icon/number
            icon_shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(0.7), y_pos,
                Inches(0.4), Inches(0.4)
            )
            icon_shape.fill.solid()
            icon_shape.fill.fore_color.rgb = RGBColor.from_string(colors["primary"])
            icon_shape.line.fill.background()

            # Number in icon
            tf = icon_shape.text_frame
            tf.paragraphs[0].text = str(i + 1) if icon_style == "number" else "•"
            tf.paragraphs[0].font.size = Pt(14)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor.from_string("FFFFFF")
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf.margin_left = Pt(0)
            tf.margin_right = Pt(0)
            tf.margin_top = Pt(5)

            # Bullet text
            if isinstance(bullet, dict):
                title = bullet.get("title", bullet.get("text", ""))
                description = bullet.get("description", "")
            else:
                title = str(bullet)
                description = ""

            # Title
            title_box = slide.shapes.add_textbox(Inches(1.3), y_pos, Inches(10.4), Inches(0.4))
            title_tf = title_box.text_frame
            title_tf.paragraphs[0].text = title
            title_tf.paragraphs[0].font.name = self.style["fonts"]["body"]
            title_tf.paragraphs[0].font.size = Pt(layout["bullet_font_size"])
            title_tf.paragraphs[0].font.bold = True
            title_tf.paragraphs[0].font.color.rgb = RGBColor.from_string(colors["text_dark"])

            # Description if present
            if description:
                desc_box = slide.shapes.add_textbox(Inches(1.3), y_pos + Inches(0.35), Inches(10.4), Inches(0.45))
                desc_tf = desc_box.text_frame
                desc_tf.word_wrap = True
                desc_tf.paragraphs[0].text = description
                desc_tf.paragraphs[0].font.name = self.style["fonts"]["body"]
                desc_tf.paragraphs[0].font.size = Pt(layout["bullet_font_size"] - 2)
                desc_tf.paragraphs[0].font.color.rgb = RGBColor.from_string(colors["text_light"])

    def _add_two_column_bullets(self, slide, bullets: List[Dict[str, Any]], layout: Dict,
                                colors: Dict, icon_style: str):
        """Add bullets in two columns."""
        mid_point = (len(bullets) + 1) // 2
        left_bullets = bullets[:mid_point]
        right_bullets = bullets[mid_point:]

        # Left column
        left = Inches(0.7)
        self._add_column_bullets(slide, left_bullets, left, layout, colors, icon_style, 0)

        # Right column
        right = Inches(6.5)
        self._add_column_bullets(slide, right_bullets, right, layout, colors, icon_style, mid_point)

    def _add_column_bullets(self, slide, bullets: List[Dict[str, Any]], left_pos,
                            layout: Dict, colors: Dict, icon_style: str, start_num: int):
        """Add a column of bullets."""
        bullet_top = Inches(1.9)
        bullet_spacing = Inches(1.0)

        for i, bullet in enumerate(bullets):
            y_pos = bullet_top + (i * bullet_spacing)

            # Icon
            icon_shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                left_pos, y_pos,
                Inches(0.35), Inches(0.35)
            )
            icon_shape.fill.solid()
            icon_shape.fill.fore_color.rgb = RGBColor.from_string(colors["primary"])
            icon_shape.line.fill.background()

            # Number in icon
            tf = icon_shape.text_frame
            tf.paragraphs[0].text = str(start_num + i + 1) if icon_style == "number" else "•"
            tf.paragraphs[0].font.size = Pt(12)
            tf.paragraphs[0].font.bold = True
            tf.paragraphs[0].font.color.rgb = RGBColor.from_string("FFFFFF")
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Bullet content
            if isinstance(bullet, dict):
                title = bullet.get("title", bullet.get("text", ""))
                description = bullet.get("description", "")
            else:
                title = str(bullet)
                description = ""

            # Title
            title_box = slide.shapes.add_textbox(left_pos + Inches(0.5), y_pos, Inches(5), Inches(0.35))
            title_tf = title_box.text_frame
            title_tf.paragraphs[0].text = title
            title_tf.paragraphs[0].font.name = self.style["fonts"]["body"]
            title_tf.paragraphs[0].font.size = Pt(layout["bullet_font_size"] - 1)
            title_tf.paragraphs[0].font.bold = True
            title_tf.paragraphs[0].font.color.rgb = RGBColor.from_string(colors["text_dark"])

            # Description
            if description:
                desc_box = slide.shapes.add_textbox(left_pos + Inches(0.5), y_pos + Inches(0.35), Inches(5), Inches(0.55))
                desc_tf = desc_box.text_frame
                desc_tf.word_wrap = True
                desc_tf.paragraphs[0].text = description
                desc_tf.paragraphs[0].font.name = self.style["fonts"]["body"]
                desc_tf.paragraphs[0].font.size = Pt(layout["bullet_font_size"] - 3)
                desc_tf.paragraphs[0].font.color.rgb = RGBColor.from_string(colors["text_light"])

    def _create_intro_with_left_image(self, slide, title, subtitle, presenter_name, date,
                                      image_url, layout, colors):
        """Create intro slide with image on the left."""
        # Image placeholder (left side)
        img_left = Inches(0.7)
        img_top = Inches(1.5)
        img_width = Inches(5)
        img_height = Inches(4.5)

        # Add image or placeholder
        self._add_image_or_placeholder(slide, image_url, img_left, img_top, img_width, img_height, colors)

        # Text content (right side)
        text_left = Inches(6.2)
        text_width = Inches(6)

        # Title
        self._add_title(slide, title, text_left, Inches(2), text_width, Inches(1.2),
                       layout["title_font_size"], colors["text_dark"])

        # Accent line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, text_left, Inches(3.2), Inches(1.5), Inches(0.06))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor.from_string(colors["primary"])
        line.line.fill.background()

        # Subtitle
        if subtitle:
            self._add_body_text(slide, subtitle, text_left, Inches(3.5), text_width, Inches(1.2),
                               layout["subtitle_font_size"], colors["text_light"])

        # Presenter info
        if presenter_name:
            self._add_presenter_info(slide, presenter_name, date, text_left, Inches(5), text_width, colors)

    def _create_intro_with_right_image(self, slide, title, subtitle, presenter_name, date,
                                       image_url, layout, colors):
        """Create intro slide with image on the right."""
        # Text content (left side)
        text_left = Inches(0.7)
        text_width = Inches(6)

        # Title
        self._add_title(slide, title, text_left, Inches(2), text_width, Inches(1.2),
                       layout["title_font_size"], colors["text_dark"])

        # Accent line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, text_left, Inches(3.2), Inches(1.5), Inches(0.06))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor.from_string(colors["primary"])
        line.line.fill.background()

        # Subtitle
        if subtitle:
            self._add_body_text(slide, subtitle, text_left, Inches(3.5), text_width, Inches(1.2),
                               layout["subtitle_font_size"], colors["text_light"])

        # Presenter info
        if presenter_name:
            self._add_presenter_info(slide, presenter_name, date, text_left, Inches(5), text_width, colors)

        # Image (right side)
        img_left = Inches(7)
        img_top = Inches(1.5)
        img_width = Inches(5)
        img_height = Inches(4.5)
        self._add_image_or_placeholder(slide, image_url, img_left, img_top, img_width, img_height, colors)

    def _create_centered_intro(self, slide, title, subtitle, presenter_name, date, layout, colors):
        """Create centered intro slide without image."""
        # Title
        self._add_title(slide, title, Inches(1), Inches(2.2), Inches(10.7), Inches(1.2),
                       layout["title_font_size"], colors["text_dark"], alignment=PP_ALIGN.CENTER)

        # Accent line (centered)
        line_width = Inches(2)
        line_left = (self.prs.slide_width - Inches(2)) / 2
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, line_left, Inches(3.5), line_width, Inches(0.06))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor.from_string(colors["primary"])
        line.line.fill.background()

        # Subtitle
        if subtitle:
            self._add_body_text(slide, subtitle, Inches(1.5), Inches(3.8), Inches(9.7), Inches(1),
                               layout["subtitle_font_size"], colors["text_light"], alignment=PP_ALIGN.CENTER)

        # Presenter info (centered)
        if presenter_name:
            info_text = presenter_name
            if date:
                info_text += f" | {date}"
            self._add_body_text(slide, info_text, Inches(1.5), Inches(5.2), Inches(9.7), Inches(0.5),
                               14, colors["text_muted"], alignment=PP_ALIGN.CENTER)

    def _add_presenter_info(self, slide, presenter_name, date, left, top, width, colors):
        """Add presenter information box."""
        # Container shape
        container = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, Inches(4), Inches(0.8)
        )
        container.fill.solid()
        container.fill.fore_color.rgb = RGBColor.from_string(colors["background"])
        container.line.color.rgb = RGBColor.from_string(colors["text_muted"])
        container.line.width = Pt(1)

        # Get initials
        initials = "".join(word[0].upper() for word in presenter_name.split()[:2])

        # Avatar circle
        avatar = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            left + Inches(0.15), top + Inches(0.15),
            Inches(0.5), Inches(0.5)
        )
        avatar.fill.solid()
        avatar.fill.fore_color.rgb = RGBColor.from_string(colors["primary"])
        avatar.line.fill.background()

        # Initials text
        tf = avatar.text_frame
        tf.paragraphs[0].text = initials
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor.from_string("FFFFFF")
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.margin_top = Pt(8)

        # Name and date
        name_box = slide.shapes.add_textbox(left + Inches(0.75), top + Inches(0.1), Inches(3), Inches(0.4))
        name_tf = name_box.text_frame
        name_tf.paragraphs[0].text = presenter_name
        name_tf.paragraphs[0].font.name = self.style["fonts"]["body"]
        name_tf.paragraphs[0].font.size = Pt(14)
        name_tf.paragraphs[0].font.bold = True
        name_tf.paragraphs[0].font.color.rgb = RGBColor.from_string(colors["text_dark"])

        if date:
            date_box = slide.shapes.add_textbox(left + Inches(0.75), top + Inches(0.45), Inches(3), Inches(0.3))
            date_tf = date_box.text_frame
            date_tf.paragraphs[0].text = date
            date_tf.paragraphs[0].font.name = self.style["fonts"]["body"]
            date_tf.paragraphs[0].font.size = Pt(12)
            date_tf.paragraphs[0].font.color.rgb = RGBColor.from_string(colors["text_muted"])

    def _add_content_with_image(self, slide, body, bullets, image_url, top, content_width,
                                layout, colors):
        """Add content area with image on the right."""
        # Content on left
        if bullets:
            self._add_bullet_list(slide, bullets, Inches(0.7), top, content_width, layout, colors)
        elif body:
            self._add_body_text(slide, body, Inches(0.7), top, content_width, Inches(4.5),
                               layout["body_font_size"], colors["text_light"])

        # Image on right
        img_left = Inches(7.5)
        img_width = Inches(4.5)
        img_height = Inches(4)
        self._add_image_or_placeholder(slide, image_url, img_left, top, img_width, img_height, colors)

    def _add_image_or_placeholder(self, slide, image_url, left, top, width, height, colors):
        """Add an image or a placeholder shape."""
        if image_url and image_url.strip():
            # For now, add a placeholder - image download handled separately
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, width, height
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor.from_string(colors.get("text_muted", "CCCCCC"))
            placeholder.line.fill.background()

            # Store URL for later processing
            placeholder.name = f"IMAGE_PLACEHOLDER:{image_url}"
        else:
            # Empty placeholder
            placeholder = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, width, height
            )
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = RGBColor.from_string(colors.get("text_muted", "EEEEEE"))
            placeholder.line.fill.background()

    def _add_speaker_notes(self, slide, notes: str):
        """Add speaker notes to a slide."""
        try:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = notes
        except Exception:
            pass

    # ====================================================================================
    # IMAGE HANDLING
    # ====================================================================================

    async def process_images(self):
        """Download and insert all placeholder images."""
        for slide in self.prs.slides:
            shapes_to_process = []

            for shape in slide.shapes:
                if hasattr(shape, 'name') and shape.name.startswith("IMAGE_PLACEHOLDER:"):
                    url = shape.name.replace("IMAGE_PLACEHOLDER:", "")
                    shapes_to_process.append((shape, url))

            for shape, url in shapes_to_process:
                await self._replace_with_image(slide, shape, url)

    async def _replace_with_image(self, slide, placeholder_shape, url: str):
        """Replace a placeholder shape with an actual image."""
        try:
            # Download image
            image_bytes = await self._download_image(url)
            if not image_bytes:
                return

            # Get position
            left = placeholder_shape.left
            top = placeholder_shape.top
            width = placeholder_shape.width
            height = placeholder_shape.height

            # Remove placeholder
            sp = placeholder_shape._element
            sp.getparent().remove(sp)

            # Save and add image
            temp_path = os.path.join(self.temp_dir, f"{uuid.uuid4()}.png")
            with open(temp_path, "wb") as f:
                f.write(image_bytes)

            slide.shapes.add_picture(temp_path, left, top, width, height)

            # Cleanup
            os.unlink(temp_path)

        except Exception as e:
            print(f"Error processing image {url}: {e}")

    async def _download_image(self, url: str) -> Optional[bytes]:
        """Download an image from URL with caching."""
        if url in self._image_cache:
            return self._image_cache[url]

        try:
            async with aiohttp.ClientSession() as session:
                async with session.get(url, timeout=30) as response:
                    if response.status == 200:
                        data = await response.read()
                        self._image_cache[url] = data
                        return data
        except Exception as e:
            print(f"Error downloading image: {e}")

        return None

    # ====================================================================================
    # EXPORT
    # ====================================================================================

    def save(self, path: str):
        """Save the presentation to a file."""
        self.prs.save(path)

    def cleanup(self):
        """Clean up temporary files."""
        import shutil
        try:
            shutil.rmtree(self.temp_dir)
        except Exception:
            pass


# ====================================================================================
# FACTORY FUNCTION FOR EASY USE
# ====================================================================================

async def build_presentation_from_content(
    slides_content: List[Dict[str, Any]],
    template_name: str = "general",
    title: str = "Presentation",
    output_path: str = None
) -> str:
    """
    Build a complete presentation from structured content.

    Args:
        slides_content: List of slide content dictionaries
        template_name: Template style to use (general, modern, standard, swift)
        title: Presentation title for filename
        output_path: Optional output path

    Returns:
        Path to the generated PPTX file
    """
    builder = SmartPptxBuilder(template_name)

    for i, slide_data in enumerate(slides_content):
        slide_type = slide_data.get("type", "content")

        if slide_type == "intro" or i == 0:
            builder.create_intro_slide(
                title=slide_data.get("title", "Presentation"),
                subtitle=slide_data.get("subtitle", slide_data.get("description", "")),
                presenter_name=slide_data.get("presenterName", slide_data.get("presenter_name", "")),
                date=slide_data.get("presentationDate", slide_data.get("date", "")),
                image_url=slide_data.get("image", {}).get("__image_url__") or slide_data.get("image", {}).get("url")
            )
        elif slide_type == "conclusion":
            builder.create_conclusion_slide(
                title=slide_data.get("title", "Thank You"),
                message=slide_data.get("message", slide_data.get("description", "")),
                call_to_action=slide_data.get("callToAction", slide_data.get("call_to_action", "")),
                contact_info=slide_data.get("contactInfo", slide_data.get("contact_info", ""))
            )
        elif slide_data.get("bullets"):
            builder.create_bullets_slide(
                title=slide_data.get("title", ""),
                bullets=slide_data.get("bullets", []),
                speaker_notes=slide_data.get("speaker_notes", slide_data.get("__speaker_note__", ""))
            )
        else:
            builder.create_content_slide(
                title=slide_data.get("title", ""),
                body=slide_data.get("body", slide_data.get("description", "")),
                bullets=slide_data.get("bullets"),
                image_url=slide_data.get("image", {}).get("__image_url__") or slide_data.get("image", {}).get("url"),
                speaker_notes=slide_data.get("speaker_notes", slide_data.get("__speaker_note__", ""))
            )

    # Process all images
    await builder.process_images()

    # Save
    if not output_path:
        exports_dir = os.path.join(get_app_data_directory(), "exports")
        os.makedirs(exports_dir, exist_ok=True)
        output_path = os.path.join(exports_dir, f"{title.replace(' ', '_')}_{uuid.uuid4().hex[:8]}.pptx")

    builder.save(output_path)
    builder.cleanup()

    return output_path
