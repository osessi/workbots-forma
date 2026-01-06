"""
Initialize Built-in Templates as PPTX files for Smart Templates

This module creates professional PPTX template files for the built-in templates
(general, modern, standard, swift) so they can be used with the Smart Templates
system which preserves original design during export.

The templates are automatically created on first startup if they don't exist.
You can replace them with better designed PPTX files via the admin interface.
"""

import os
import uuid
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from sqlmodel import Session, select

from services.database import engine
from models.sql.template import PptxTemplateModel
from utils.datetime_utils import get_current_utc_datetime
from utils.get_env import get_app_data_directory_env as get_app_data_directory


# Template color schemes matching the React templates
TEMPLATE_SCHEMES = {
    "general": {
        "name": "Général",
        "description": "Template polyvalent pour tous types de présentations avec layouts flexibles",
        "primary": "9333EA",      # Purple
        "title_color": "111827",  # Dark gray
        "body_color": "4B5563",   # Gray
        "muted_color": "6B7280",  # Muted gray
        "bg_color": "FFFFFF",     # White
        "font_heading": "Inter",
        "font_body": "Inter",
    },
    "modern": {
        "name": "Moderne",
        "description": "Design épuré avec accents bleus, idéal pour pitch decks",
        "primary": "1E4CD9",      # Blue
        "title_color": "1E4CD9",  # Blue
        "body_color": "374151",   # Dark gray
        "muted_color": "6B7280",  # Muted gray
        "bg_color": "FFFFFF",     # White
        "font_heading": "Montserrat",
        "font_body": "Montserrat",
    },
    "standard": {
        "name": "Standard",
        "description": "Format classique et professionnel avec mise en page structurée",
        "primary": "1B8C2D",      # Green
        "title_color": "111827",  # Dark gray
        "body_color": "4B5563",   # Gray
        "muted_color": "6B7280",  # Muted gray
        "bg_color": "FFFFFF",     # White
        "font_heading": "Playfair Display",
        "font_body": "Inter",
    },
    "swift": {
        "name": "Swift",
        "description": "Style minimaliste et dynamique avec couleurs sky blue",
        "primary": "0EA5E9",      # Sky blue
        "title_color": "111827",  # Dark gray
        "body_color": "4B5563",   # Gray
        "muted_color": "6B7280",  # Muted gray
        "bg_color": "FFFFFF",     # White
        "font_heading": "Inter",
        "font_body": "Inter",
    },
}


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to RGBColor"""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


def create_title_slide(prs: Presentation, scheme: dict):
    """Create a title/intro slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add accent bar at top
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(0.1)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = hex_to_rgb(scheme["primary"])
    accent_bar.line.fill.background()

    # Add title placeholder
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(11), Inches(1.5)
    )
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_p.text = "Titre de la présentation"
    title_p.font.name = scheme["font_heading"]
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = hex_to_rgb(scheme["title_color"])
    title_p.alignment = PP_ALIGN.CENTER

    # Add subtitle placeholder
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.2),
        Inches(11), Inches(0.8)
    )
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.word_wrap = True
    subtitle_p = subtitle_tf.paragraphs[0]
    subtitle_p.text = "Sous-titre ou description courte"
    subtitle_p.font.name = scheme["font_body"]
    subtitle_p.font.size = Pt(20)
    subtitle_p.font.color.rgb = hex_to_rgb(scheme["body_color"])
    subtitle_p.alignment = PP_ALIGN.CENTER

    # Add presenter info at bottom
    presenter_box = slide.shapes.add_textbox(
        Inches(1), Inches(6),
        Inches(11), Inches(0.5)
    )
    presenter_tf = presenter_box.text_frame
    presenter_p = presenter_tf.paragraphs[0]
    presenter_p.text = "Présentateur • Date"
    presenter_p.font.name = scheme["font_body"]
    presenter_p.font.size = Pt(14)
    presenter_p.font.color.rgb = hex_to_rgb(scheme["muted_color"])
    presenter_p.alignment = PP_ALIGN.CENTER

    return slide


def create_content_slide(prs: Presentation, scheme: dict):
    """Create a content slide with title and body"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add accent bar at top
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(0.08)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = hex_to_rgb(scheme["primary"])
    accent_bar.line.fill.background()

    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.5),
        Inches(11.4), Inches(1)
    )
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_p.text = "Titre de la slide"
    title_p.font.name = scheme["font_heading"]
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = hex_to_rgb(scheme["title_color"])

    # Add accent line under title
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.5),
        Inches(1.2), Inches(0.05)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = hex_to_rgb(scheme["primary"])
    accent_line.line.fill.background()

    # Add body text area
    body_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.8),
        Inches(11.4), Inches(4.5)
    )
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    body_p = body_tf.paragraphs[0]
    body_p.text = "Contenu de la slide. Ajoutez votre texte ici."
    body_p.font.name = scheme["font_body"]
    body_p.font.size = Pt(18)
    body_p.font.color.rgb = hex_to_rgb(scheme["body_color"])

    return slide


def create_bullets_slide(prs: Presentation, scheme: dict):
    """Create a slide with bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add accent bar at top
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(0.08)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = hex_to_rgb(scheme["primary"])
    accent_bar.line.fill.background()

    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.5),
        Inches(11.4), Inches(1)
    )
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_p.text = "Points clés"
    title_p.font.name = scheme["font_heading"]
    title_p.font.size = Pt(36)
    title_p.font.bold = True
    title_p.font.color.rgb = hex_to_rgb(scheme["title_color"])

    # Add accent line under title
    accent_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(1.5),
        Inches(1.2), Inches(0.05)
    )
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = hex_to_rgb(scheme["primary"])
    accent_line.line.fill.background()

    # Add bullet points
    bullets_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.9),
        Inches(11.4), Inches(4.5)
    )
    bullets_tf = bullets_box.text_frame
    bullets_tf.word_wrap = True

    bullet_items = [
        "Premier point important",
        "Deuxième point à retenir",
        "Troisième élément clé",
        "Quatrième aspect essentiel"
    ]

    for i, item in enumerate(bullet_items):
        if i == 0:
            p = bullets_tf.paragraphs[0]
        else:
            p = bullets_tf.add_paragraph()

        p.text = f"• {item}"
        p.font.name = scheme["font_body"]
        p.font.size = Pt(20)
        p.font.color.rgb = hex_to_rgb(scheme["body_color"])
        p.space_before = Pt(12)
        p.space_after = Pt(6)

    return slide


def create_image_content_slide(prs: Presentation, scheme: dict):
    """Create a slide with image placeholder and text"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add accent bar at top
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(0.08)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = hex_to_rgb(scheme["primary"])
    accent_bar.line.fill.background()

    # Add title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.5),
        Inches(5.5), Inches(1)
    )
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_p.text = "Titre avec image"
    title_p.font.name = scheme["font_heading"]
    title_p.font.size = Pt(32)
    title_p.font.bold = True
    title_p.font.color.rgb = hex_to_rgb(scheme["title_color"])

    # Add description
    desc_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.6),
        Inches(5.5), Inches(4)
    )
    desc_tf = desc_box.text_frame
    desc_tf.word_wrap = True
    desc_p = desc_tf.paragraphs[0]
    desc_p.text = "Description du contenu. Vous pouvez ajouter plusieurs paragraphes de texte ici pour accompagner l'image."
    desc_p.font.name = scheme["font_body"]
    desc_p.font.size = Pt(16)
    desc_p.font.color.rgb = hex_to_rgb(scheme["body_color"])

    # Add image placeholder (gray rectangle)
    img_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.8), Inches(1),
        Inches(5.5), Inches(5)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = RGBColor(229, 231, 235)  # Gray-200
    img_placeholder.line.color.rgb = RGBColor(209, 213, 219)  # Gray-300

    # Add "Image" text in placeholder
    img_text = slide.shapes.add_textbox(
        Inches(6.8), Inches(3.2),
        Inches(5.5), Inches(0.6)
    )
    img_tf = img_text.text_frame
    img_p = img_tf.paragraphs[0]
    img_p.text = "Image"
    img_p.font.name = scheme["font_body"]
    img_p.font.size = Pt(18)
    img_p.font.color.rgb = RGBColor(156, 163, 175)  # Gray-400
    img_p.alignment = PP_ALIGN.CENTER

    return slide


def create_closing_slide(prs: Presentation, scheme: dict):
    """Create a closing/thank you slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add large accent rectangle at bottom
    accent_rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.5),
        prs.slide_width, Inches(2)
    )
    accent_rect.fill.solid()
    accent_rect.fill.fore_color.rgb = hex_to_rgb(scheme["primary"])
    accent_rect.line.fill.background()

    # Add thank you text
    thanks_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(11), Inches(1.5)
    )
    thanks_tf = thanks_box.text_frame
    thanks_p = thanks_tf.paragraphs[0]
    thanks_p.text = "Merci !"
    thanks_p.font.name = scheme["font_heading"]
    thanks_p.font.size = Pt(54)
    thanks_p.font.bold = True
    thanks_p.font.color.rgb = hex_to_rgb(scheme["title_color"])
    thanks_p.alignment = PP_ALIGN.CENTER

    # Add contact info
    contact_box = slide.shapes.add_textbox(
        Inches(1), Inches(6),
        Inches(11), Inches(0.8)
    )
    contact_tf = contact_box.text_frame
    contact_p = contact_tf.paragraphs[0]
    contact_p.text = "contact@example.com • www.example.com"
    contact_p.font.name = scheme["font_body"]
    contact_p.font.size = Pt(16)
    contact_p.font.color.rgb = RGBColor(255, 255, 255)
    contact_p.alignment = PP_ALIGN.CENTER

    return slide


def create_template_pptx(template_id: str, scheme: dict, output_dir: str) -> str:
    """Create a complete PPTX template file"""
    prs = Presentation()

    # Set slide size to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Create slides
    create_title_slide(prs, scheme)
    create_content_slide(prs, scheme)
    create_bullets_slide(prs, scheme)
    create_image_content_slide(prs, scheme)
    create_closing_slide(prs, scheme)

    # Save
    output_path = os.path.join(output_dir, f"{template_id}.pptx")
    prs.save(output_path)

    return output_path


def init_builtin_templates_if_needed():
    """
    Initialize built-in templates if they don't already exist.
    Called at API startup.
    """
    # Create templates directory
    templates_dir = os.path.join(get_app_data_directory(), "smart_templates")
    os.makedirs(templates_dir, exist_ok=True)

    print("Checking built-in PPTX templates...")

    with Session(engine) as session:
        for template_id, scheme in TEMPLATE_SCHEMES.items():
            # Check if already exists
            existing = session.exec(
                select(PptxTemplateModel).where(
                    PptxTemplateModel.name == scheme["name"],
                    PptxTemplateModel.is_active == True
                )
            ).first()

            if existing:
                # Check if file still exists
                if existing.file_path and os.path.exists(existing.file_path):
                    continue
                else:
                    # File missing, recreate it
                    print(f"  - Recreating {scheme['name']} (file was missing)")
                    pptx_path = create_template_pptx(template_id, scheme, templates_dir)
                    existing.file_path = pptx_path
                    existing.updated_at = get_current_utc_datetime()
                    session.add(existing)
                    session.commit()
                    continue

            # Create new template
            print(f"  - Creating {scheme['name']}...")
            pptx_path = create_template_pptx(template_id, scheme, templates_dir)

            # Get file size
            file_size = os.path.getsize(pptx_path)

            # Create database record
            template = PptxTemplateModel(
                id=uuid.uuid4(),
                name=scheme["name"],
                description=scheme["description"],
                category="builtin",
                file_path=pptx_path,
                file_size=file_size,
                slide_count=5,
                placeholder_mapping=[],
                slide_layouts={
                    "title": 0,
                    "content": 1,
                    "bullets": 2,
                    "image": 3,
                    "closing": 4
                },
                fonts=[scheme["font_heading"], scheme["font_body"]],
                theme_colors={
                    "primary": scheme["primary"],
                    "title": scheme["title_color"],
                    "body": scheme["body_color"],
                    "muted": scheme["muted_color"],
                    "background": scheme["bg_color"]
                },
                is_active=True,
                is_system=True,
                created_at=get_current_utc_datetime(),
                updated_at=get_current_utc_datetime()
            )

            session.add(template)
            session.commit()
            print(f"    Created: {pptx_path}")

    print("Built-in templates check complete.")
